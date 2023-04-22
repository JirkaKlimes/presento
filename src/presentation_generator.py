import collections.abc
from pptx import Presentation
from pathlib import Path
import os
import random

class PPTXGenerator:

    LAYOUT_TITLE = 0
    LAYOUT_TITLE_AND_CONTENT = 1

    def __init__(self):
        self._load_themes()

    def _load_themes(self):
        self.themes = os.listdir("./data/themes")

    def generate(self, parsed_presentation, file_name):

        theme = random.choice(self.themes)

        prs = Presentation(f"./data/themes/{theme}")

        for psl in parsed_presentation:

            match psl["layout"]:
                case "Title Slide":
                    layout = prs.slide_layouts[self.LAYOUT_TITLE]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = psl["heading"]
                    if psl["subheading"]:
                        slide.placeholders[1].text = psl["subheading"]


                case "Title and Content":
                    layout = prs.slide_layouts[self.LAYOUT_TITLE_AND_CONTENT]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = psl["heading"]

                    bullet_shape = slide.shapes.placeholders[1]
                    tf = bullet_shape.text_frame
                    for i, text in enumerate(psl['points']):
                        if i == 0:
                            tf.text = text
                        else:
                            p = tf.add_paragraph()
                            p.text = text


        Path('./output').mkdir(parents=True, exist_ok=True)
        prs.save(f'./output/{file_name}.pptx')


if __name__ == "__main__":
    import json
    pres_gen = PPTXGenerator()

    with open("./debug/parsed_presentation.json", encoding='utf-8') as f:
        parsed = json.load(f)

    pres_gen.generate(parsed, 'output')